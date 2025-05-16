from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Создание презентации
    prs = Presentation()
    
    # Настройка размеров слайдов (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Стиль текста
    title_font_size = Pt(44)
    subtitle_font_size = Pt(28)
    bullet_font_size = Pt(24)
    
    # Цвета
    title_color = RGBColor(0, 32, 96)  # Темно-синий
    text_color = RGBColor(0, 0, 0)     # Черный
    accent_color = RGBColor(0, 112, 192)  # Синий
    
    # === Слайд 1: Титульный слайд ===
    slide_layout = prs.slide_layouts[0]  # Титульный слайд
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Мультимедийный ПК для домашнего кинотеатра и VR"
    title.text_frame.paragraphs[0].font.size = title_font_size
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Проект высокопроизводительного ПК для 8K-кинотеатра и виртуальной реальности"
    subtitle.text_frame.paragraphs[0].font.size = subtitle_fontont_size
    
    # === Слайд 2: Цель проекта ===
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Цель проекта"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()  # Очистка стандартного текста
    
    p = tf.add_paragraph()
    p.text = "Спроектировать ПК для:"
    p.font.size = bullet_font_size
    
    p = tf.add_paragraph()
    p.text = "▪ Просмотра 8K-фильмов с Dolby Atmos"
    p.level = 1
    p.font.size = bullet_font_size
    
    p = tf.add_paragraph()
    p.text = "▪ Игр в виртуальной реальности (Meta Quest 3, Valve Index)"
    p.level = 1
    p.font.size = bullet_font_size
    
    p = tf.add_paragraph()
    p.text = "▪ Бесшумной работы в гостиной (<30 дБ)"
    p.level = 1
    p.font.size = bullet_font_size
    
    # === Слайд 3: Требования к железу ===
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Требования к комплектующим"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    components = [
        ("Процессор", "12+ ядер (Intel i7-14700 / AMD Ryzen 9 7900)"),
        ("Материнская плата", "ATX, PCIe 5.0, Wi-Fi 7 (Z790 / X670)"),
        ("Оперативная память", "32 ГБ DDR5 (5200+ МГц)"),
        ("Накопитель", "NVMe SSD 2 ТБ (PCIe 5.0)"),
        ("Видеокарта", "NVIDIA RTX 4080 / AMD RX 7900 XT (12-16 ГБ VRAM)"),
        ("Охлаждение", "Кастомное жидкостное"),
        ("Блок питания", "850-1000 Вт (80 PLUS Gold, модульный)"),
        ("Сеть", "Ethernet 2.5 Гбит/с + Wi-Fi 7")
    ]
    
    for component, desc in components:
        p = tf.add_paragraph()
        p.text = f"▪ {component}: {desc}"
        p.font.size = bullet_font_size
        p.level = 0
    
    # === Слайд 4: Бюджет и особенности ===
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Бюджет и ключевые особенности"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    features = [
        "Поддержка HDMI 2.1 (8K @ 120 Гц)",
        "Тихая работа (<30 дБ) благодаря жидкостному охлаждению",
        "Поддержка Dolby Atmos и DTS:X",
        "Быстрая загрузка игр и 8K-контента",
        "Совместимость с VR-шлемами (Wi-Fi 7 для беспроводного режима)",
        "Бюджет: $2000 – $3000"
    ]
    
    for feature in features:
        p = tf.add_paragraph()
        p.text = f"▪ {feature}"
        p.font.size = bullet_font_size
        p.level = 0
    
    # === Слайд 5: Сравнение процессоров ===
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Варианты процессоров"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    # Создаем таблицу
    rows, cols = 3, 3
    left = Inches(1.5)
    top = Inches(2)
    width = Inches(8)
    height = Inches(3)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Заголовки таблицы
    table.cell(0, 0).text = "Модель"
    table.cell(0, 1).text = "Ядра/Потоки"
    table.cell(0, 2).text = "Частота (Boost)"
    
    # Данные
    cpus = [
        ("Intel Core i7-14700", "12/20", "5.4 ГГц"),
        ("AMD Ryzen 9 7900", "12/24", "5.6 ГГц")
    ]
    
    for i, cpu in enumerate(cpus, start=1):
        table.cell(i, 0).text = cpu[0]
        table.cell(i, 1).text = cpu[1]
        table.cell(i, 2).text = cpu[2]
    
    # Стиль таблицы
    for cell in table.iter_cells():
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.alignment = PP_ALIGN.CENTER
    
    # === Слайд 6: Сравнение видеокарт ===
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Варианты видеокарт"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    # Создаем таблицу
    rows, cols = 3, 4
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(10)
    height = Inches(3)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Заголовки
    headers = ["Модель", "VRAM", "Память", "HDMI 2.1"]
    for col, header in enumerate(headers):
        table.cell(0, col).text = header
    
    # Данные
    gpus = [
        ("NVIDIA RTX 4080", "16 ГБ GDDR6X", "256 бит", "Да"),
        ("AMD RX 7900 XT", "20 ГБ GDDR6", "320 бит", "Да")
    ]
    
    for i, gpu in enumerate(gpus, start=1):
        for j, value in enumerate(gpu):
            table.cell(i, j).text = value
    
    # Стиль таблицы
    for cell in table.iter_cells():
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.alignment = PP_ALIGN.CENTER
    
    # === Слайд 7: Итоговая схема ===
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Итоговая конфигурация"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    final_config = [
        "Процессор: AMD Ryzen 9 7900 (12/24, 5.6 ГГц)",
        "Видеокарта: NVIDIA RTX 4080 (16 ГБ GDDR6X)",
        "Память: 32 ГБ DDR5-5600",
        "Накопитель: 2 ТБ NVMe SSD (PCIe 5.0)",
        "Охлаждение: Кастомная СЖО с низкооборотистыми вентиляторами",
        "Блок питания: 1000 Вт 80+ Platinum (модульный)",
        "Корпус: Fractal Design Define 7 (шумопоглощение)"
    ]
    
    for item in final_config:
        p = tf.add_paragraph()
        p.text = f"▪ {item}"
        p.font.size = bullet_font_size
        p.level = 0
    
    # Сохранение презентации
    prs.save("Мультимедийный_PC_для_домашнего_кинотеатра_и_VR.pptx")
    print("Презентация успешно создана!")

if __name__ == "__main__":
    create_presentation()
